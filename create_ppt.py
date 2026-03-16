"""
Create PowerPoint presentation from ppt2.md using python-pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def parse_markdown_table(lines: list[str]) -> tuple[list[list[str]], int] | None:
    """Parse a markdown table. Returns (rows, num_lines_consumed) or None."""
    if not lines or "|" not in lines[0]:
        return None
    rows = []
    i = 0
    for i, line in enumerate(lines):
        line = line.strip()
        if not line or "|" not in line:
            break
        # Skip separator row (|---|---|)
        if re.match(r"^\|[\s\-:]+\|", line):
            continue
        cells = [c.strip() for c in line.split("|") if c.strip() or line.count("|") > 1]
        if cells:
            rows.append(cells)
    return (rows, i) if rows else None


def parse_code_block(lines: list[str]) -> tuple[str, int] | None:
    """Parse a fenced code block. Returns (content, num_lines_consumed) or None."""
    if not lines or not lines[0].strip().startswith("```"):
        return None
    content = []
    i = 1
    while i < len(lines):
        if lines[i].strip().startswith("```"):
            i += 1
            break
        content.append(lines[i])
        i += 1
    return ("\n".join(content), i) if content else ("", i)


def strip_markdown(text: str) -> str:
    """Remove markdown formatting for plain text."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text.strip()


def parse_slide(content: str) -> dict:
    """Parse a slide's markdown content into structured data."""
    lines = content.strip().split("\n")
    result = {
        "title": "",
        "subtitle": "",
        "body": [],
        "tables": [],
        "code_blocks": [],
    }
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # ## Slide N: Title
        if stripped.startswith("## "):
            result["title"] = re.sub(r"^##\s*(?:Slide\s*\d+:\s*)?", "", stripped).strip()
            result["title"] = strip_markdown(result["title"])
            i += 1
            continue

        # # Main title (large heading)
        if stripped.startswith("# ") and not stripped.startswith("## "):
            txt = re.sub(r"^#+\s*", "", stripped).strip()
            result["subtitle"] = strip_markdown(txt)
            i += 1
            continue

        # Code block
        code = parse_code_block(lines[i:])
        if code:
            result["code_blocks"].append(code[0])
            i += code[1]
            continue

        # Table
        table = parse_markdown_table(lines[i:])
        if table:
            result["tables"].append(table[0])
            i += table[1]
            continue

        # Bullet or paragraph
        if stripped.startswith("- "):
            result["body"].append(("bullet", strip_markdown(stripped[2:])))
        elif stripped:
            result["body"].append(("para", strip_markdown(stripped)))
        i += 1

    return result


def add_text_to_frame(tf, text: str, bold: bool = False, size: int = 12):
    """Add text to a text frame with optional formatting."""
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold


def add_paragraph_to_frame(tf, text: str, level: int = 0, size: int = 11):
    """Add a new paragraph to text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)


def create_slide(prs: Presentation, parsed: dict, slide_num: int):
    """Create a single slide from parsed content."""
    # Use blank layout for flexibility
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.8)

    # Title
    if parsed["title"]:
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = parsed["title"]
        p.font.size = Pt(24)
        p.font.bold = True
        top += Inches(0.6)

    # Subtitle (e.g. main title on slide 1)
    if parsed["subtitle"]:
        tx = slide.shapes.add_textbox(left, top, width, Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = parsed["subtitle"]
        p.font.size = Pt(20)
        p.font.bold = True
        top += Inches(1.0)

    # Body text (bullets and paragraphs)
    body_lines = []
    for kind, text in parsed["body"]:
        if kind == "bullet":
            body_lines.append(("• " + text, 0))
        else:
            body_lines.append((text, -1))  # -1 = paragraph

    if body_lines:
        tx = slide.shapes.add_textbox(left, top, width, Inches(2.5))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, (text, level) in enumerate(body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(11)
            if level >= 0:
                p.level = level
        top += Inches(2.2)

    # Tables
    for table_rows in parsed["tables"]:
        if not table_rows:
            continue
        rows, cols = len(table_rows), max(len(r) for r in table_rows)
        tbl_height = min(Inches(0.4) * rows, Inches(2.5))
        tbl = slide.shapes.add_table(rows, cols, left, top, width, tbl_height).table
        for r, row in enumerate(table_rows):
            for c, cell_text in enumerate(row):
                if c < cols:
                    tbl.cell(r, c).text = cell_text[:50]  # truncate long text
        top += tbl_height + Inches(0.3)

    # Code blocks
    for code in parsed["code_blocks"]:
        code_lines = code.strip().split("\n")
        block_height = min(Inches(0.25) * len(code_lines) + Inches(0.2), Inches(3))
        tx = slide.shapes.add_textbox(left, top, width, block_height)
        tf = tx.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = code.strip()[:800]  # limit length
        p.font.size = Pt(9)
        p.font.name = "Consolas"
        top += block_height + Inches(0.2)

    return slide


def main():
    print("Loading ppt2.md...")
    md_path = Path(__file__).parent / "ppt2.md"
    out_path = Path(__file__).parent / "ppt2.pptx"

    content = md_path.read_text(encoding="utf-8")
    raw_slides = re.split(r"\n---\n", content)

    # Filter out empty and header-only blocks
    slides = []
    for block in raw_slides:
        block = block.strip()
        if not block:
            continue
        # Skip the initial header (first 3 lines before first ---)
        slides.append(block)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print(f"Creating {len(slides)} slides...", flush=True)
    for i, block in enumerate(slides):
        parsed = parse_slide(block)
        create_slide(prs, parsed, i + 1)
        if (i + 1) % 5 == 0 or i == 0:
            print(f"  Slide {i + 1}/{len(slides)}...", flush=True)

    print("Saving...", flush=True)
    prs.save(out_path)
    print(f"Created: {out_path}")


if __name__ == "__main__":
    main()
