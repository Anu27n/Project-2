"""Quick test of python-pptx."""
import sys
print("Starting...", flush=True)
from pptx import Presentation
from pptx.util import Inches
print("Creating presentation...", flush=True)
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.paragraphs[0].text = "Test"
print("Saving...", flush=True)
prs.save("test_out.pptx")
print("Done!", flush=True)
